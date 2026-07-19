from __future__ import annotations

import io

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Emu

from crewmeal.search_enhancement.geometry_facts import (
    _extract_gantt,
    geometry_facts_by_slide,
)

TABLE_LEFT = 1_000_000
TABLE_TOP = 1_000_000
LABEL_W = 2_000_000
WEEK_W = 400_000
ROW_H = 300_000
BAR_H = 180_000


def _week_left(week_index: int) -> int:
    # week_index is 1-based (W1..W4)
    return TABLE_LEFT + LABEL_W + (week_index - 1) * WEEK_W


def _build_gantt_deck() -> bytes:
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    rows, cols = 4, 5  # month row, week row, Task A, Task B
    frame = slide.shapes.add_table(
        rows, cols, Emu(TABLE_LEFT), Emu(TABLE_TOP),
        Emu(LABEL_W + 4 * WEEK_W), Emu(4 * ROW_H),
    )
    table = frame.table
    table.columns[0].width = Emu(LABEL_W)
    for ci in range(1, 5):
        table.columns[ci].width = Emu(WEEK_W)
    for ri in range(rows):
        table.rows[ri].height = Emu(ROW_H)

    # month header: 1월 covers only W1, 2월 covers W2-W4 (non-uniform span)
    table.cell(0, 1).text = "1월"
    table.cell(0, 2).text = "2월"
    table.cell(0, 2).merge(table.cell(0, 4))
    # week header
    for ci, wk in zip(range(1, 5), ["W1", "W2", "W3", "W4"]):
        table.cell(1, ci).text = wk
    table.cell(2, 0).text = "Task A"
    table.cell(3, 0).text = "Task B"

    # Task A bar covers W1..W2
    a = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Emu(_week_left(1)), Emu(TABLE_TOP + 2 * ROW_H + 50_000),
        Emu(2 * WEEK_W), Emu(BAR_H),
    )
    a.name = "bar-a"
    # Task B bar covers W4 only
    b = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Emu(_week_left(4)), Emu(TABLE_TOP + 3 * ROW_H + 50_000),
        Emu(WEEK_W), Emu(BAR_H),
    )
    b.name = "bar-b"
    # Milestone triangle centered on W3
    tri = slide.shapes.add_shape(
        MSO_SHAPE.ISOSCELES_TRIANGLE,
        Emu(_week_left(3) + WEEK_W // 2 - 72_000), Emu(TABLE_TOP + 2 * ROW_H + 50_000),
        Emu(144_000), Emu(BAR_H),
    )
    tri.name = "ms"
    label = slide.shapes.add_textbox(
        Emu(_week_left(3)), Emu(TABLE_TOP + 2 * ROW_H - 120_000), Emu(WEEK_W), Emu(120_000)
    )
    label.text_frame.text = "GO-LIVE"

    buffer = io.BytesIO()
    prs.save(buffer)
    return buffer.getvalue()


def test_gantt_month_spans_and_task_weeks_are_deterministic():
    deck = _build_gantt_deck()
    prs = Presentation(io.BytesIO(deck))
    gantt = _extract_gantt(prs.slides[0])

    assert gantt is not None
    assert gantt["month_spans"] == {"1월": "W1", "2월": "W2-W4"}

    schedule = {item["task"]: (item["start"], item["end"]) for item in gantt["schedule"]}
    assert schedule["Task A"] == ("1월 W1", "2월 W2")
    assert schedule["Task B"] == ("2월 W4", "2월 W4")

    milestones = {m["label"]: m["week"] for m in gantt["milestones"]}
    assert milestones.get("GO-LIVE") == "2월 W3"


def test_geometry_facts_by_slide_emits_ground_truth_block():
    facts = geometry_facts_by_slide(_build_gantt_deck())
    assert set(facts) == {1}
    text = facts[1]
    assert "1월=W1" in text and "2월=W2-W4" in text
    assert "Task A: 1월 W1 ~ 2월 W2" in text
    assert "Task B: 2월 W4" in text
    assert "GO-LIVE @ 2월 W3" in text


def test_non_gantt_slide_yields_no_facts():
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    box = slide.shapes.add_textbox(Emu(1_000_000), Emu(1_000_000), Emu(3_000_000), Emu(500_000))
    box.text_frame.text = "Just a title, no schedule"
    buffer = io.BytesIO()
    prs.save(buffer)
    assert geometry_facts_by_slide(buffer.getvalue()) == {}
