from __future__ import annotations

import io

from pptx import Presentation
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.oxml.ns import qn
from pptx.util import Emu

from crewmeal.search_enhancement.geometry_facts import (
    _LEVEL_TOLERANCE,
    _extract_diagram,
    _orient,
    geometry_facts_by_slide,
)

EMU_IN = 914_400
BOX_W, BOX_H = int(1.6 * EMU_IN), int(0.7 * EMU_IN)

# An org tree is just one kind of node-link diagram: hierarchy lives ONLY in
# geometry (vertical level + connectors), box text is names only. With no
# arrowheads the extractor must infer top-down direction from layout.
BOXES = {
    "대표이사": (4.5, 0.6),
    "개발본부": (3.0, 2.0),
    "사업본부": (6.0, 2.0),
    "플랫폼팀": (2.0, 3.4),
    "데이터팀": (3.6, 3.4),
    "영업팀": (6.0, 3.4),
}
EDGES = [
    ("대표이사", "개발본부"),
    ("대표이사", "사업본부"),
    ("개발본부", "플랫폼팀"),
    ("개발본부", "데이터팀"),
    ("사업본부", "영업팀"),
]


def _build_org_deck(connect: bool = True) -> bytes:
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    shapes = {}
    for name, (cx_in, top_in) in BOXES.items():
        left = int(cx_in * EMU_IN) - BOX_W // 2
        top = int(top_in * EMU_IN)
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Emu(left), Emu(top), Emu(BOX_W), Emu(BOX_H))
        box.text_frame.text = name
        shapes[name] = box
    for parent, child in EDGES:
        p, c = shapes[parent], shapes[child]
        conn = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT,
            Emu(p.left + p.width // 2), Emu(p.top + p.height),
            Emu(c.left + c.width // 2), Emu(c.top),
        )
        if connect:
            conn.begin_connect(p, 2)
            conn.end_connect(c, 0)
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _add_tail_arrow(conn) -> None:
    """Inject <a:tailEnd type="triangle"/> so the connector points start -> end."""
    spPr = conn._element.find(qn("p:spPr"))
    ln = spPr.find(qn("a:ln"))
    if ln is None:
        ln = spPr.makeelement(qn("a:ln"), {})
        spPr.append(ln)
    ln.append(ln.makeelement(qn("a:tailEnd"), {"type": "triangle"}))


def _build_flow_deck(arrow: bool = True, same_level: bool = True) -> bytes:
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    names = ["수집", "처리", "적재"]
    boxes = {}
    for i, n in enumerate(names):
        left = int((1.0 + i * 2.2) * EMU_IN)
        top = int((3.0 if same_level else 2.0 + i) * EMU_IN)
        b = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(left), Emu(top), Emu(BOX_W), Emu(BOX_H))
        b.text_frame.text = n
        boxes[n] = b
    for a, c in [("수집", "처리"), ("처리", "적재")]:
        p, q = boxes[a], boxes[c]
        conn = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT,
            Emu(p.left + p.width), Emu(p.top + p.height // 2),
            Emu(q.left), Emu(q.top + q.height // 2),
        )
        conn.begin_connect(p, 3)
        conn.end_connect(q, 1)
        if arrow:
            _add_tail_arrow(conn)
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _build_two_flow_deck() -> bytes:
    """One slide with two physically separate arrow chains and no connector
    between them, so the extractor must report them as distinct flows."""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    def _chain(names: list[str], top_in: float) -> None:
        boxes = {}
        for i, n in enumerate(names):
            left = int((1.0 + i * 2.2) * EMU_IN)
            top = int(top_in * EMU_IN)
            b = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, Emu(left), Emu(top), Emu(BOX_W), Emu(BOX_H)
            )
            b.text_frame.text = n
            boxes[n] = b
        for a, c in zip(names, names[1:]):
            p, q = boxes[a], boxes[c]
            conn = slide.shapes.add_connector(
                MSO_CONNECTOR.STRAIGHT,
                Emu(p.left + p.width), Emu(p.top + p.height // 2),
                Emu(q.left), Emu(q.top + q.height // 2),
            )
            conn.begin_connect(p, 3)
            conn.end_connect(q, 1)
            _add_tail_arrow(conn)

    _chain(["수집", "처리", "적재"], 1.5)
    _chain(["요청", "응답"], 4.5)
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


# --- org tree (a vertical-layout diagram) --------------------------------------


def test_org_tree_edges_from_explicit_connections():
    prs = Presentation(io.BytesIO(_build_org_deck(connect=True)))
    diagram = _extract_diagram(prs.slides[0])
    assert diagram is not None
    assert set(diagram["nodes"]) == set(BOXES)
    got = {(e["source"], e["target"]) for e in diagram["edges"]}
    assert got == set(EDGES)
    assert all(e["directed"] for e in diagram["edges"])


def test_org_tree_edges_from_geometry_fallback_when_unconnected():
    prs = Presentation(io.BytesIO(_build_org_deck(connect=False)))
    diagram = _extract_diagram(prs.slides[0])
    assert diagram is not None
    got = {(e["source"], e["target"]) for e in diagram["edges"]}
    assert got == set(EDGES)


def test_geometry_facts_by_slide_emits_diagram_edges():
    facts = geometry_facts_by_slide(_build_org_deck(connect=True))
    assert set(facts) == {1}
    text = facts[1]
    assert "Diagram relationships" in text
    assert "대표이사 -> 개발본부" in text
    assert "사업본부 -> 영업팀" in text


def test_disconnected_flows_render_as_separate_groups():
    facts = geometry_facts_by_slide(_build_two_flow_deck())
    assert set(facts) == {1}
    text = facts[1]
    assert "SEPARATE" in text
    assert "Flow 1" in text
    assert "Flow 2" in text
    # each chain stays intact and is not fused with the other
    assert "수집 -> 처리" in text
    assert "처리 -> 적재" in text
    assert "요청 -> 응답" in text


# --- flowchart (direction from arrowheads, not layout) -------------------------


def test_flow_arrows_set_direction_even_on_one_level():
    diagram = _extract_diagram(Presentation(io.BytesIO(_build_flow_deck(arrow=True, same_level=True))).slides[0])
    assert diagram is not None
    edges = {(e["source"], e["target"]): e["directed"] for e in diagram["edges"]}
    assert edges.get(("수집", "처리")) is True
    assert edges.get(("처리", "적재")) is True


def test_plain_links_are_undirected_on_one_level():
    diagram = _extract_diagram(Presentation(io.BytesIO(_build_flow_deck(arrow=False, same_level=True))).slides[0])
    assert diagram is not None
    assert diagram["edges"]
    assert all(not e["directed"] for e in diagram["edges"])


# --- direction resolution unit tests -------------------------------------------


def test_orient_tail_arrow_points_start_to_end():
    a = {"id": 1, "text": "A", "cy": 100}
    b = {"id": 2, "text": "B", "cy": 100}
    edge = _orient(a, b, None, "triangle")
    assert (edge["source"], edge["target"], edge["directed"]) == ("A", "B", True)


def test_orient_head_arrow_reverses_direction():
    a = {"id": 1, "text": "A", "cy": 100}
    b = {"id": 2, "text": "B", "cy": 100}
    edge = _orient(a, b, "triangle", "none")
    assert (edge["source"], edge["target"], edge["directed"]) == ("B", "A", True)


def test_orient_double_arrow_is_undirected():
    a = {"id": 1, "text": "A", "cy": 100}
    b = {"id": 2, "text": "B", "cy": 100}
    assert _orient(a, b, "triangle", "triangle")["directed"] is False


def test_orient_layout_falls_back_to_top_down():
    top = {"id": 1, "text": "T", "cy": 100}
    bottom = {"id": 2, "text": "B", "cy": 100 + _LEVEL_TOLERANCE}
    edge = _orient(bottom, top, None, None)
    assert (edge["source"], edge["target"], edge["directed"]) == ("T", "B", True)


def test_orient_same_level_no_arrow_is_undirected():
    a = {"id": 1, "text": "A", "cy": 100}
    b = {"id": 2, "text": "B", "cy": 150}
    assert _orient(a, b, None, None)["directed"] is False


# --- guard ---------------------------------------------------------------------


def test_diagram_is_noop_without_connectors():
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    for i in range(3):
        b = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Emu(1_000_000 + i * 2_000_000), Emu(1_000_000), Emu(BOX_W), Emu(BOX_H)
        )
        b.text_frame.text = f"box{i}"
    buf = io.BytesIO()
    prs.save(buf)
    assert _extract_diagram(Presentation(io.BytesIO(buf.getvalue())).slides[0]) is None
