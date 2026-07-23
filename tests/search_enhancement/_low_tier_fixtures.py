"""python-pptx deck builders for the low-tier extractor tests.

The raw-XML ``_pptx_fixtures`` package targets the OOXML *semantic* parser, which
reads slide XML directly. The low-tier extractor instead opens the deck with
``python-pptx``, which needs a complete package (slide master, layouts, theme).
Building fixtures with ``python-pptx`` itself is the simplest way to get a valid
one, so these helpers wrap the common shape-authoring calls.

The ``tests`` tree has no ``__init__.py`` packages, so this module is imported by
bare name (pytest's prepend import mode puts the test directory on ``sys.path``).
The underscore prefix keeps pytest from collecting it as a test module.
"""

from __future__ import annotations

import io
from collections.abc import Sequence

from PIL import Image
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches, Pt

# Layout indices in python-pptx's default template.
_TITLE_ONLY = 5
_BLANK = 6


def new_presentation() -> Presentation:
    return Presentation()


def add_title_slide(prs: Presentation):
    """Add a slide that has a real title placeholder."""

    return prs.slides.add_slide(prs.slide_layouts[_TITLE_ONLY])


def add_blank_slide(prs: Presentation):
    """Add a slide with no placeholders (title must be inferred)."""

    return prs.slides.add_slide(prs.slide_layouts[_BLANK])


def set_title(slide, text: str) -> None:
    slide.shapes.title.text = text


def add_textbox(
    slide, *lines: str, top_in: float = 2.0, left_in: float = 1.0, font_pt: int | None = None
):
    box = slide.shapes.add_textbox(
        Inches(left_in), Inches(top_in), Inches(4), Inches(1)
    )
    frame = box.text_frame
    for index, line in enumerate(lines):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.text = line
        if font_pt is not None:
            paragraph.font.size = Pt(font_pt)
    return box


def add_table(slide, rows: Sequence[Sequence[str]], *, top_in: float = 2.0):
    row_count = len(rows)
    col_count = max(len(row) for row in rows)
    frame = slide.shapes.add_table(
        row_count, col_count, Inches(1), Inches(top_in), Inches(6), Inches(2)
    )
    table = frame.table
    for r, row in enumerate(rows):
        for c in range(col_count):
            table.cell(r, c).text = row[c] if c < len(row) else ""
    return table


def add_chart(
    slide,
    categories: Sequence[str],
    series_name: str,
    values: Sequence[float],
    *,
    title: str | None = None,
):
    data = CategoryChartData()
    data.categories = list(categories)
    data.add_series(series_name, tuple(values))
    graphic_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(1),
        Inches(2),
        Inches(6),
        Inches(4),
        data,
    )
    chart = graphic_frame.chart
    if title is not None:
        chart.has_title = True
        chart.chart_title.text_frame.text = title
    return chart


def make_png(color: tuple[int, int, int] = (180, 40, 40), size: tuple[int, int] = (80, 60)) -> bytes:
    buffer = io.BytesIO()
    Image.new("RGB", size, color).save(buffer, format="PNG")
    return buffer.getvalue()


def add_picture(slide, png: bytes | None = None, *, top_in: float = 2.0):
    blob = png if png is not None else make_png()
    return slide.shapes.add_picture(
        io.BytesIO(blob), Inches(1), Inches(top_in), Inches(2), Inches(1.5)
    )


def add_group_textbox(slide, *lines: str, top_in: float = 3.0):
    group = slide.shapes.add_group_shape()
    box = group.shapes.add_textbox(
        Inches(1), Inches(top_in), Inches(4), Inches(1)
    )
    frame = box.text_frame
    for index, line in enumerate(lines):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.text = line
    return group


def set_notes(slide, text: str) -> None:
    slide.notes_slide.notes_text_frame.text = text


def render(prs: Presentation) -> bytes:
    buffer = io.BytesIO()
    prs.save(buffer)
    return buffer.getvalue()
