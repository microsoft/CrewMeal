"""Deterministic slide geometry facts extracted from PPTX (OOXML) coordinates.

The vision model cannot reliably read exact positions from a downsampled slide
image (Gantt bar week columns, milestone placement). Those positions are encoded
exactly in the file: table column widths, row heights, header merges and the
absolute EMU coordinates of every bar/marker autoshape. This module reconstructs
that structure deterministically so it can be injected into analysis as
ground-truth evidence, instead of asking the model to eyeball pixels.

Currently implemented: Gantt / timeline schedule slides (table with a month
header row + a ``W#`` week header row and free-floating bar/milestone shapes).
The extractor is a no-op (returns ``None``) for slides that do not match, so it
is always safe to run over an entire deck.
"""

from __future__ import annotations

import io
from collections import Counter, defaultdict
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.oxml.ns import qn

# EMU tolerance used when deciding whether an autoshape is a schedule bar or a
# marker, and whether a bar overlaps a week column enough to count.
_HEIGHT_TOLERANCE = 30_000
_MIN_BAR_WIDTH = 150_000
_COLUMN_OVERLAP_RATIO = 0.30

# Node-link diagram reconstruction (org charts, flowcharts, architecture and
# network diagrams). Minimum text nodes / connectors / resolved edges for a
# slide to count as a diagram; the vertical gap used to infer a top-down
# direction when a connector has no arrowhead; and the maximum distance an
# unconnected connector endpoint may sit from a node box to still bind to it.
_MIN_DIAGRAM_NODES = 3
_MIN_DIAGRAM_CONNECTORS = 2
_MIN_DIAGRAM_EDGES = 2
_LEVEL_TOLERANCE = 200_000
_MAX_FALLBACK_DISTANCE = 400_000


def geometry_facts_by_slide(source: Path | bytes) -> dict[int, str]:
    """Return ``{slide_number: evidence_text}`` for slides with derivable geometry.

    Slide numbers are 1-based to match the rest of the pipeline. Slides without
    recognisable geometry are omitted. Never raises: any parsing problem yields
    an empty result so manifest/analysis flow is unaffected.
    """
    try:
        prs = _open(source)
    except Exception:
        return {}

    facts: dict[int, str] = {}
    for idx, slide in enumerate(prs.slides):
        parts: list[str] = []
        try:
            gantt = _extract_gantt(slide)
        except Exception:
            gantt = None
        if gantt and gantt["schedule"]:
            parts.append(_render_gantt_text(gantt))
        try:
            diagram = _extract_diagram(slide)
        except Exception:
            diagram = None
        if diagram and diagram["edges"]:
            parts.append(_render_diagram_text(diagram))
        if parts:
            facts[idx + 1] = "\n\n".join(parts)
    return facts


def _open(source: Path | bytes) -> Any:
    if isinstance(source, (bytes, bytearray)):
        return Presentation(io.BytesIO(bytes(source)))
    return Presentation(str(source))


def _overlap(a0: float, a1: float, b0: float, b1: float) -> float:
    return max(0.0, min(a1, b1) - max(a0, b0))


def _is_triangle(shape: Any) -> bool:
    try:
        name = (shape.auto_shape_type.name or "").upper()
        if "TRIANGLE" in name:
            return True
    except Exception:
        pass
    raw = (shape.name or "")
    return "삼각형" in raw or "TRIANGLE" in raw.upper()


def _dominant_bar_height(slide: Any, gx0: int, gx1: int, body_top: int) -> int:
    heights: Counter[int] = Counter()
    for shape in slide.shapes:
        if shape.shape_type != MSO_SHAPE_TYPE.AUTO_SHAPE or _is_triangle(shape):
            continue
        try:
            left, top, width, height = shape.left, shape.top, shape.width, shape.height
        except Exception:
            continue
        if left is None or top + height < body_top:
            continue
        center = left + width / 2
        if gx0 - 200_000 <= center <= gx1 + 200_000 and width >= _MIN_BAR_WIDTH and width > height:
            heights[round(height, -3)] += 1
    return heights.most_common(1)[0][0] if heights else 180_000


def _month_spans(week_cols: dict[int, int], month_of_col: dict[int, str]) -> dict[str, str]:
    spans: dict[str, list[int]] = {}
    for ci in sorted(week_cols):
        spans.setdefault(month_of_col[ci], []).append(week_cols[ci])
    return {
        month: (f"W{min(weeks)}-W{max(weeks)}" if len(weeks) > 1 else f"W{weeks[0]}")
        for month, weeks in spans.items()
    }


def _extract_gantt(slide: Any) -> dict[str, Any] | None:
    tbl_shape = next((s for s in slide.shapes if getattr(s, "has_table", False)), None)
    if tbl_shape is None:
        return None
    table = tbl_shape.table
    cols = list(table.columns)
    rows = list(table.rows)
    if len(cols) < 5 or len(rows) < 3:
        return None

    cx = [tbl_shape.left]
    for col in cols:
        cx.append(cx[-1] + col.width)
    ry = [tbl_shape.top]
    for row in rows:
        ry.append(ry[-1] + row.height)
    ncols = len(cols)

    def cell_text(ri: int, ci: int) -> str:
        try:
            return table.cell(ri, ci).text.replace("\n", " ").strip()
        except Exception:
            return ""

    # locate the week-header row (>=4 cells shaped like W<number>)
    week_row = None
    for ri in range(min(4, len(rows))):
        hits = sum(
            1
            for ci in range(ncols)
            if cell_text(ri, ci).upper().startswith("W") and cell_text(ri, ci)[1:].isdigit()
        )
        if hits >= 4:
            week_row = ri
            break
    if week_row is None:
        return None
    month_row = max(0, week_row - 1)

    week_cols: dict[int, int] = {}
    for ci in range(ncols):
        txt = cell_text(week_row, ci).upper()
        if txt.startswith("W") and txt[1:].isdigit():
            week_cols[ci] = int(txt[1:])
    if len(week_cols) < 4:
        return None
    first_week_col = min(week_cols)

    # month per week column, filling forward across merged month header cells
    month_of_col: dict[int, str] = {}
    current = ""
    for ci in sorted(week_cols):
        label = cell_text(month_row, ci)
        if label:
            current = label
        month_of_col[ci] = current

    def week_label(ci: int) -> str:
        month = month_of_col.get(ci, "")
        return f"{month} W{week_cols[ci]}".strip()

    body_rows = []
    for ri in range(week_row + 1, len(rows)):
        labels = [cell_text(ri, ci) for ci in range(first_week_col)]
        name = max(labels, key=len) if labels else ""
        body_rows.append((ri, ry[ri], ry[ri + 1], name, cell_text(ri, 0)))

    grid_x0, grid_x1 = cx[first_week_col], cx[-1]
    body_top = ry[week_row + 1]
    bar_h = _dominant_bar_height(slide, grid_x0, grid_x1, body_top)

    bars: list[tuple[int, int, int, int]] = []
    marks: list[tuple[int, int, int, int]] = []
    for shape in slide.shapes:
        if shape.shape_type != MSO_SHAPE_TYPE.AUTO_SHAPE:
            continue
        try:
            left, top, width, height = shape.left, shape.top, shape.width, shape.height
        except Exception:
            continue
        if left is None:
            continue
        center = left + width / 2
        if not (grid_x0 - 200_000 <= center <= grid_x1 + 200_000):
            continue
        if top + height < body_top:
            continue
        if not _is_triangle(shape) and abs(height - bar_h) <= _HEIGHT_TOLERANCE and width >= _MIN_BAR_WIDTH:
            bars.append((left, top, width, height))
        else:
            marks.append((left, top, width, height))

    def weeks_for_xrange(x0: int, x1: int) -> list[int]:
        return [
            ci
            for ci in sorted(week_cols)
            if _overlap(x0, x1, cx[ci], cx[ci + 1]) > _COLUMN_OVERLAP_RATIO * cols[ci].width
        ]

    # one bar (widest) per task row
    per_row: dict[int, tuple[int, int, int, int]] = {}
    for left, top, width, height in bars:
        center_y = top + height / 2
        key = next((ri for (ri, y0, y1, _n, _c) in body_rows if y0 <= center_y < y1), None)
        if key is None:
            continue
        if key not in per_row or width > per_row[key][2]:
            per_row[key] = (left, top, width, height)

    schedule = []
    for key in sorted(per_row):
        left, top, width, _h = per_row[key]
        name, category = next((n, c) for (ri, _0, _1, n, c) in body_rows if ri == key)
        hit = weeks_for_xrange(left, left + width)
        if not hit:
            continue
        start, end = min(hit), max(hit)
        schedule.append(
            {
                "task": name,
                "category": category,
                "start": week_label(start),
                "end": week_label(end),
            }
        )

    textboxes = [
        (s.left, s.top, s.width, s.height, s.text_frame.text.replace("\n", " ").strip())
        for s in slide.shapes
        if s.has_text_frame
        and s.shape_type == MSO_SHAPE_TYPE.TEXT_BOX
        and s.text_frame.text.strip()
    ]
    milestones = []
    for left, top, width, height in marks:
        center_x = left + width / 2
        center_y = top + height / 2
        ci = min(week_cols, key=lambda c: abs((cx[c] + cx[c + 1]) / 2 - center_x))
        label = ""
        best = float("inf")
        for lx, lt, lw, lh, txt in textboxes:
            dist = ((lx + lw / 2) - center_x) ** 2 + ((lt + lh / 2) - center_y) ** 2
            if dist < best:
                best, label = dist, txt
        if label:
            milestones.append({"label": label, "week": week_label(ci)})

    return {
        "month_spans": _month_spans(week_cols, month_of_col),
        "schedule": schedule,
        "milestones": milestones,
    }


def _render_gantt_text(gantt: dict[str, Any]) -> str:
    axis = ", ".join(f"{m}={span}" for m, span in gantt["month_spans"].items())
    lines = [
        "Gantt/timeline schedule reconstructed from exact table + bar coordinates.",
        f"Month axis (each month spans a variable number of week columns): {axis}",
        "Task schedule (first week column ~ last week column the coloured bar covers):",
    ]
    for item in gantt["schedule"]:
        prefix = f"{item['category']} / " if item["category"] and item["category"] != item["task"] else ""
        span = item["start"] if item["start"] == item["end"] else f"{item['start']} ~ {item['end']}"
        lines.append(f"- {prefix}{item['task']}: {span}")
    if gantt["milestones"]:
        marks = "; ".join(f"{m['label']} @ {m['week']}" for m in gantt["milestones"])
        lines.append(f"Milestones (single week column each): {marks}")
    return "\n".join(lines)


def _child_transform(group: Any, parent_t: Any, parent_scale: tuple[float, float]) -> tuple[Any, tuple[float, float]]:
    """Return ``(transform, scale)`` mapping a group's child coordinate space to
    slide EMU. Grouped shapes report ``.left``/``.top`` and connector endpoints in
    the group's child space (``a:chOff``/``a:chExt``), not slide space, so every
    grouped shape must be projected through this before its geometry is comparable
    to top-level shapes. Scales compose for nested groups.
    """
    try:
        x, y, cx, cy = group.left, group.top, group.width, group.height
        xfrm = group._element.find(qn("p:grpSpPr") + "/" + qn("a:xfrm"))
        choff = xfrm.find(qn("a:chOff")) if xfrm is not None else None
        chext = xfrm.find(qn("a:chExt")) if xfrm is not None else None
        if x is None or choff is None or chext is None:
            return parent_t, parent_scale
        chx, chy = int(choff.get("x")), int(choff.get("y"))
        chcx, chcy = int(chext.get("cx")), int(chext.get("cy"))
    except Exception:
        return parent_t, parent_scale
    sx = (cx / chcx) if chcx else 1.0
    sy = (cy / chcy) if chcy else 1.0

    def transform(px: float, py: float) -> tuple[float, float]:
        return parent_t(x + (px - chx) * sx, y + (py - chy) * sy)

    return transform, (parent_scale[0] * sx, parent_scale[1] * sy)


def _collect_diagram_shapes(
    shapes: Any,
    transform: Any,
    scale: tuple[float, float],
    nodes: dict[int, dict[str, Any]],
    connectors: list[tuple[Any, Any]],
) -> None:
    """Recurse into groups, projecting every leaf into absolute slide EMU.

    Text-bearing shapes become nodes (keyed by shape id, with an absolute bbox);
    connectors are collected with the transform that maps their endpoints to slide
    space. Explicit ``stCxn``/``endCxn`` references are shape-id based and thus
    coordinate-independent; the transform only matters for the geometric fallback.
    """
    for shape in shapes:
        try:
            stype = shape.shape_type
        except Exception:
            continue
        if stype == MSO_SHAPE_TYPE.GROUP:
            child_t, child_scale = _child_transform(shape, transform, scale)
            _collect_diagram_shapes(shape.shapes, child_t, child_scale, nodes, connectors)
            continue
        if stype == MSO_SHAPE_TYPE.LINE:
            connectors.append((shape, transform))
            continue
        try:
            if not shape.has_text_frame:
                continue
            text = " ".join(shape.text_frame.text.split())
            if not text or shape.left is None:
                continue
            ax, ay = transform(shape.left, shape.top)
            aw, ah = shape.width * scale[0], shape.height * scale[1]
        except Exception:
            continue
        nodes[shape.shape_id] = {
            "id": shape.shape_id,
            "text": text,
            "left": ax,
            "top": ay,
            "width": aw,
            "height": ah,
            "cy": ay + ah / 2,
        }


def _cxn_shape_id(conn: Any, tag: str) -> int | None:
    el = conn._element.find(qn("p:nvCxnSpPr") + "/" + qn("p:cNvCxnSpPr") + "/" + qn(tag))
    if el is None:
        return None
    raw = el.get("id")
    return int(raw) if raw is not None else None


def _arrow_ends(conn: Any) -> tuple[str | None, str | None]:
    """Return ``(headEnd, tailEnd)`` arrow types. ``headEnd`` sits at the connector
    start point, ``tailEnd`` at the end point; ``None``/``"none"`` means no arrow."""
    ln = conn._element.find(qn("p:spPr") + "/" + qn("a:ln"))
    if ln is None:
        return (None, None)
    head = ln.find(qn("a:headEnd"))
    tail = ln.find(qn("a:tailEnd"))
    return (
        head.get("type") if head is not None else None,
        tail.get("type") if tail is not None else None,
    )


def _node_distance(px: float, py: float, node: dict[str, Any]) -> float:
    dx = max(node["left"] - px, 0, px - (node["left"] + node["width"]))
    dy = max(node["top"] - py, 0, py - (node["top"] + node["height"]))
    return (dx * dx + dy * dy) ** 0.5


def _resolve_node(
    shape_id: int | None,
    px: float | None,
    py: float | None,
    transform: Any,
    nodes: dict[int, dict[str, Any]],
    node_list: list[dict[str, Any]],
) -> dict[str, Any] | None:
    if shape_id is not None and shape_id in nodes:
        return nodes[shape_id]
    if px is None or py is None or not node_list:
        return None
    ax, ay = transform(px, py)
    best = min(node_list, key=lambda n: _node_distance(ax, ay, n))
    return best if _node_distance(ax, ay, best) <= _MAX_FALLBACK_DISTANCE else None


def _orient(
    start: dict[str, Any], end: dict[str, Any], head_type: str | None, tail_type: str | None
) -> dict[str, Any]:
    """Decide edge direction: arrowheads first (a tail arrow points start->end),
    then top-down layout when the two nodes sit on different levels, else the edge
    is left undirected."""
    arrow_at_start = head_type not in (None, "none")
    arrow_at_end = tail_type not in (None, "none")
    if arrow_at_end and not arrow_at_start:
        src, tgt, directed = start, end, True
    elif arrow_at_start and not arrow_at_end:
        src, tgt, directed = end, start, True
    elif arrow_at_start and arrow_at_end:
        src, tgt, directed = start, end, False
    elif abs(start["cy"] - end["cy"]) >= _LEVEL_TOLERANCE:
        src, tgt = (start, end) if start["cy"] < end["cy"] else (end, start)
        directed = True
    else:
        src, tgt, directed = start, end, False
    return {
        "source": src["text"],
        "target": tgt["text"],
        "source_id": src["id"],
        "target_id": tgt["id"],
        "directed": directed,
    }


def _extract_diagram(slide: Any) -> dict[str, Any] | None:
    """Reconstruct a node-link diagram (text boxes = nodes, connectors = edges).

    Generic across org charts, flowcharts, process flows and architecture/network
    diagrams: no diagram type is assumed. Edges bind by explicit connector
    ``stCxn``/``endCxn`` shape references when present (robust inside groups),
    otherwise by nearest node to each connector endpoint. Direction comes from
    connector arrowheads, falling back to top-down layout, else the edge is
    undirected. Returns ``None`` for slides that are not connected-box diagrams so
    it is safe to run over an entire deck.
    """
    nodes: dict[int, dict[str, Any]] = {}
    connectors: list[tuple[Any, Any]] = []
    _collect_diagram_shapes(slide.shapes, lambda x, y: (x, y), (1.0, 1.0), nodes, connectors)
    if len(nodes) < _MIN_DIAGRAM_NODES or len(connectors) < _MIN_DIAGRAM_CONNECTORS:
        return None

    node_list = list(nodes.values())
    edges: list[dict[str, Any]] = []
    seen: set[tuple[int, int, bool]] = set()
    for conn, transform in connectors:
        st_id, en_id = _cxn_shape_id(conn, "a:stCxn"), _cxn_shape_id(conn, "a:endCxn")
        try:
            bx, by, ex, ey = conn.begin_x, conn.begin_y, conn.end_x, conn.end_y
        except Exception:
            bx = by = ex = ey = None
        start = _resolve_node(st_id, bx, by, transform, nodes, node_list)
        end = _resolve_node(en_id, ex, ey, transform, nodes, node_list)
        if start is None or end is None or start["id"] == end["id"]:
            continue
        head, tail = _arrow_ends(conn)
        edge = _orient(start, end, head, tail)
        if edge["directed"]:
            key = (edge["source_id"], edge["target_id"], True)
        else:
            lo, hi = sorted((edge["source_id"], edge["target_id"]))
            key = (lo, hi, False)
        if key in seen:
            continue
        seen.add(key)
        edges.append(edge)

    if len(edges) < _MIN_DIAGRAM_EDGES:
        return None
    return {"nodes": [n["text"] for n in node_list], "edges": edges}


def _edge_components(edges: list[dict[str, Any]]) -> list[list[dict[str, Any]]]:
    """Group edges into connected components by shape-id adjacency, ordered by
    node count descending.

    Physically separate flows drawn on one slide (e.g. an as-is process beside a
    to-be process) share no connector, so they land in distinct buckets. Emitting
    them separately stops the model being told they are a single graph and merging
    one flow's tail into the other.
    """
    adjacency: dict[int, set[int]] = defaultdict(set)
    for edge in edges:
        adjacency[edge["source_id"]].add(edge["target_id"])
        adjacency[edge["target_id"]].add(edge["source_id"])

    component_of: dict[int, int] = {}
    components: list[set[int]] = []
    for node_id in adjacency:
        if node_id in component_of:
            continue
        index = len(components)
        members: set[int] = set()
        stack = [node_id]
        while stack:
            current = stack.pop()
            if current in component_of:
                continue
            component_of[current] = index
            members.add(current)
            stack.extend(adjacency[current] - set(component_of))
        components.append(members)

    buckets: list[list[dict[str, Any]]] = [[] for _ in components]
    for edge in edges:
        buckets[component_of[edge["source_id"]]].append(edge)
    order = sorted(
        range(len(components)),
        key=lambda i: (len(components[i]), -i),
        reverse=True,
    )
    return [buckets[i] for i in order]


def _render_diagram_text(diagram: dict[str, Any]) -> str:
    edges = diagram["edges"]
    header = (
        f"Diagram relationships reconstructed from shape connectors "
        f"({len(diagram['nodes'])} nodes, {len(edges)} links)."
    )
    legend = "(A -> B = arrow / top-down, A -- B = plain link)"
    components = _edge_components(edges)

    if len(components) <= 1:
        lines = [header, f"Connections {legend}:"]
        for edge in edges:
            sep = "->" if edge["directed"] else "--"
            lines.append(f"- {edge['source']} {sep} {edge['target']}")
        return "\n".join(lines)

    lines = [
        header,
        f"They form {len(components)} SEPARATE, physically-disconnected flows. "
        "Keep each flow below as its own diagram; do NOT merge nodes or edges "
        "across flows even when labels repeat.",
        legend,
    ]
    for position, bucket in enumerate(components, 1):
        member_ids = {e["source_id"] for e in bucket} | {e["target_id"] for e in bucket}
        lines.append(
            f"Flow {position} ({len(member_ids)} nodes, {len(bucket)} links):"
        )
        for edge in bucket:
            sep = "->" if edge["directed"] else "--"
            lines.append(f"- {edge['source']} {sep} {edge['target']}")
    return "\n".join(lines)
