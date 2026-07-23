"""No-Vision ("low tier") structured extraction from PowerPoint slides.

The high tier renders every slide and sends it to the Vision LLM. The low tier
skips LibreOffice and Vision entirely: it reads the slide shape tree with
``python-pptx`` and produces a :class:`SlideContent` per slide from

* free text boxes and layout placeholders (title + body),
* native tables (as :class:`ContentTable` grids),
* native charts (as :class:`ContentChart` data points), and
* text recognized from embedded raster images via a local OCR engine
  (:class:`crewmeal.search_enhancement.ocr.OcrEngine`).

Unlike :mod:`crewmeal.search_enhancement.pptx_semantic` -- whose deliberately
conservative gate keeps only unambiguous linear-text slides and defers the rest
to Vision -- this module extracts *everything it can* from every slide, because
in the low tier there is no Vision fallback. Its output is intentionally
literal: it recovers text and structure, not the interpretive summaries, flow
narration, or chart insights that only the Vision model produces.

Speaker notes are never read here, matching the Vision pipeline (notes are
carried separately on the source manifest and must not leak into slide content).
"""

from __future__ import annotations

import io
import logging
from dataclasses import dataclass

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from crewmeal.search_enhancement.models import (
    ChartDataPoint,
    ContentChart,
    ContentImage,
    ContentSection,
    ContentTable,
    SlideContent,
    SlideSchedule,
)
from crewmeal.search_enhancement.ocr import OcrEngine

logger = logging.getLogger(__name__)

# A leading number badge like "01" / "3." often outranks the real title by font
# size, so the title heuristic ignores short numeric-only tokens.
_MAX_TITLE_CHARS = 60


@dataclass(slots=True)
class _ShapeText:
    top: int
    left: int
    max_font: float
    lines: tuple[str, ...]
    is_title_ph: bool


def extract_low_tier_slides(
    data: bytes, *, ocr: OcrEngine | None = None
) -> tuple[SlideContent, ...]:
    """Build one :class:`SlideContent` per slide without rendering or Vision."""

    presentation = Presentation(io.BytesIO(data))
    slides: list[SlideContent] = []
    for index, slide in enumerate(presentation.slides, start=1):
        slides.append(_slide_content(index, slide, ocr))
    return tuple(slides)


def _slide_content(slide_number: int, slide, ocr: OcrEngine | None) -> SlideContent:
    texts: list[_ShapeText] = []
    tables: list[ContentTable] = []
    charts: list[ContentChart] = []
    image_blobs: list[bytes] = []
    vector_images = 0

    def walk(shapes) -> None:
        nonlocal vector_images
        for shape in shapes:
            shape_type = shape.shape_type
            if shape_type == MSO_SHAPE_TYPE.GROUP:
                walk(shape.shapes)
                continue
            if getattr(shape, "has_table", False):
                table = _table(shape)
                if table is not None:
                    tables.append(table)
                continue
            if getattr(shape, "has_chart", False):
                chart = _chart(shape)
                if chart is not None:
                    charts.append(chart)
                continue
            if shape_type == MSO_SHAPE_TYPE.PICTURE:
                blob = _picture_blob(shape)
                if blob is not None:
                    image_blobs.append(blob)
                else:
                    vector_images += 1
                continue
            collected = _shape_text(shape)
            if collected is not None:
                texts.append(collected)

    walk(slide.shapes)

    title = _explicit_title(slide) or _heuristic_title(texts)
    body_lines = _body_lines(texts, title)
    images = _ocr_images(image_blobs, ocr)

    sections: tuple[ContentSection, ...] = ()
    if body_lines:
        sections = (
            ContentSection(heading="본문", paragraphs=(), bullets=body_lines),
        )

    warnings: list[str] = []
    if vector_images and (ocr is None or ocr.available):
        warnings.append(
            f"{vector_images}개의 벡터 이미지는 OCR로 텍스트를 추출할 수 없습니다"
            " (고품질 티어의 Vision 분석이 필요)."
        )
    if image_blobs and (ocr is None or not ocr.available):
        warnings.append(
            f"{len(image_blobs)}개의 이미지가 있으나 OCR 엔진을 사용할 수 없어"
            " 이미지 내 텍스트를 추출하지 못했습니다."
        )

    return SlideContent(
        slide_number=slide_number,
        title=title or f"슬라이드 {slide_number}",
        summary="",
        facts=(),
        sections=sections,
        hierarchies=(),
        schedule=SlideSchedule(time_axis=(), tasks=(), milestones=()),
        flows=(),
        tables=tuple(tables),
        charts=tuple(charts),
        relationships=(),
        images=images,
        warnings=tuple(warnings),
    )


def _shape_text(shape) -> _ShapeText | None:
    if not getattr(shape, "has_text_frame", False):
        return None
    lines = tuple(
        line
        for paragraph in shape.text_frame.paragraphs
        if (line := paragraph.text.strip())
    )
    if not lines:
        return None
    return _ShapeText(
        top=_int(getattr(shape, "top", None)),
        left=_int(getattr(shape, "left", None)),
        max_font=_max_font_pt(shape),
        lines=lines,
        is_title_ph=_is_title_placeholder(shape),
    )


def _explicit_title(slide) -> str:
    try:
        title_shape = slide.shapes.title
    except (ValueError, AttributeError):
        return ""
    if title_shape is None or not title_shape.has_text_frame:
        return ""
    return " ".join(
        line.strip()
        for line in title_shape.text_frame.text.splitlines()
        if line.strip()
    ).strip()


def _heuristic_title(texts: list[_ShapeText]) -> str:
    """Pick a title when no title placeholder exists.

    Prefer an explicit title placeholder; otherwise the topmost short, non-badge
    text line, using font size only to break ties between shapes at the same
    height. This mirrors how a reader scans a slide top-down.
    """

    explicit = [t for t in texts if t.is_title_ph]
    pool = explicit or texts
    candidates = [
        t for t in pool if any(_is_title_candidate(line) for line in t.lines)
    ]
    if not candidates:
        return ""
    best = min(candidates, key=lambda t: (t.top, -t.max_font, t.left))
    for line in best.lines:
        if _is_title_candidate(line):
            return line
    return ""


def _is_title_candidate(line: str) -> bool:
    stripped = line.strip()
    if not stripped or len(stripped) > _MAX_TITLE_CHARS:
        return False
    # A bare number/badge ("01", "3.", "#2") is not a title.
    return any(ch.isalpha() for ch in stripped)


def _body_lines(texts: list[_ShapeText], title: str) -> tuple[str, ...]:
    ordered = sorted(texts, key=lambda t: (t.top, t.left))
    seen: set[str] = set()
    lines: list[str] = []
    title_norm = title.strip().casefold()
    for shape_text in ordered:
        for line in shape_text.lines:
            key = line.strip().casefold()
            if not key or key == title_norm or key in seen:
                continue
            seen.add(key)
            lines.append(line)
    return tuple(lines)


def _table(shape) -> ContentTable | None:
    try:
        table = shape.table
        grid = [
            tuple((cell.text or "").strip() for cell in row.cells)
            for row in table.rows
        ]
    except Exception as exc:  # noqa: BLE001 - malformed table: skip, don't fail
        logger.debug("Skipping unreadable table: %s", exc)
        return None
    grid = [row for row in grid if any(cell for cell in row)]
    if not grid:
        return None
    headers = grid[0]
    rows = tuple(row for row in grid[1:])
    width = len(headers)
    rows = tuple(row for row in rows if len(row) == width)
    return ContentTable(title="", headers=headers, rows=rows, key_facts=())


def _chart(shape) -> ContentChart | None:
    try:
        chart = shape.chart
        categories = [str(cat) for cat in chart.plots[0].categories]
        data_points: list[ChartDataPoint] = []
        for series in chart.series:
            series_name = str(getattr(series, "name", "") or "")
            values = list(series.values)
            for position, value in enumerate(values):
                label = (
                    categories[position]
                    if position < len(categories)
                    else str(position + 1)
                )
                data_points.append(
                    ChartDataPoint(
                        series=series_name,
                        label=label,
                        value=_fmt_value(value),
                    )
                )
    except Exception as exc:  # noqa: BLE001 - unreadable chart: skip
        logger.debug("Skipping unreadable chart: %s", exc)
        return None
    if not data_points:
        return None
    title = ""
    try:
        if chart.has_title and chart.chart_title.has_text_frame:
            title = chart.chart_title.text_frame.text.strip()
    except Exception:  # noqa: BLE001
        title = ""
    return ContentChart(title=title, data_points=tuple(data_points), insights=())


def _picture_blob(shape) -> bytes | None:
    """Return the raster bytes of a picture, or ``None`` for vector/failed images."""

    try:
        image = shape.image
    except Exception:  # noqa: BLE001 - SVG/EMF/WMF/linked: no raster blob
        return None
    content_type = (getattr(image, "content_type", "") or "").lower()
    if "svg" in content_type or "emf" in content_type or "wmf" in content_type:
        return None
    try:
        return image.blob
    except Exception:  # noqa: BLE001
        return None


def _ocr_images(
    blobs: list[bytes], ocr: OcrEngine | None
) -> tuple[ContentImage, ...]:
    if not blobs or ocr is None or not ocr.available:
        return ()
    images: list[ContentImage] = []
    for blob in blobs:
        lines = ocr.read_image(blob)
        if lines:
            images.append(
                ContentImage(description="", role="ocr", visible_text=lines)
            )
    return tuple(images)


def _is_title_placeholder(shape) -> bool:
    try:
        if not shape.is_placeholder:
            return False
        return shape.placeholder_format.idx == 0
    except Exception:  # noqa: BLE001
        return False


def _max_font_pt(shape) -> float:
    best = 0.0
    try:
        for paragraph in shape.text_frame.paragraphs:
            size = paragraph.font.size
            if size is not None:
                best = max(best, size.pt)
            for run in paragraph.runs:
                if run.font.size is not None:
                    best = max(best, run.font.size.pt)
    except Exception:  # noqa: BLE001
        return best
    return best


def _int(value) -> int:
    try:
        return int(value)
    except (TypeError, ValueError):
        return 0


def _fmt_value(value) -> str:
    if isinstance(value, bool):
        return str(value)
    if isinstance(value, int):
        return str(value)
    if isinstance(value, float):
        if value.is_integer():
            return str(int(value))
        return f"{value:g}"
    return str(value).strip()
